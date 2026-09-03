"""
Document Processing Module
Handles ingestion and processing of PDF, PPTX, and other document formats
"""
import io
import re
import hashlib
from pathlib import Path
from typing import Optional
from dataclasses import dataclass, field

import pdfplumber
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import pytesseract


@dataclass
class DocumentChunk:
    """Represents a processed document chunk."""
    text: str
    metadata: dict = field(default_factory=dict)
    chunk_id: str = ""
    page_number: Optional[int] = None
    section: Optional[str] = None
    parent_chunk_id: Optional[str] = None


@dataclass
class ProcessedDocument:
    """Represents a fully processed document."""
    filename: str
    title: str
    content: str
    chunks: list[DocumentChunk]
    metadata: dict = field(default_factory=dict)
    file_hash: str = ""


class DocumentProcessor:
    """Processes various document formats into text chunks for RAG."""

    def __init__(self, docs_dir: str):
        self.docs_dir = Path(docs_dir)

    def compute_file_hash(self, filepath: Path) -> str:
        """Compute SHA256 hash of file for change detection."""
        h = hashlib.sha256()
        with open(filepath, 'rb') as f:
            for chunk in iter(lambda: f.read, b''):
                h.update(chunk)
        return h.hexdigest()[:16]

    def process_all_documents(self) -> list[ProcessedDocument]:
        """Process all supported documents in the data directory."""
        documents = []
        supported_extensions = {'.pdf', '.pptx', '.docx', '.doc', '.txt', '.png', '.jpg', '.jpeg'}

        for filepath in sorted(self.docs_dir.iterdir()):
            if filepath.suffix.lower() not in supported_extensions:
                continue
            if filepath.is_dir():
                continue

            try:
                doc = self.process_file(filepath)
                if doc:
                    documents.append(doc)
                    print(f"  ✓ Processed: {doc.filename} ({len(doc.chunks)} chunks)")
            except Exception as e:
                print(f"  ✗ Error processing {filepath.name}: {e}")

        return documents

    def process_file(self, filepath: Path) -> Optional[ProcessedDocument]:
        """Process a single file based on its extension."""
        ext = filepath.suffix.lower()
        file_hash = self.compute_file_hash(filepath)

        if ext == '.pdf':
            return self._process_pdf(filepath, file_hash)
        elif ext == '.pptx':
            return self._process_pptx(filepath, file_hash)
        elif ext in ('.docx', '.doc'):
            return self._process_docx(filepath, file_hash)
        elif ext in ('.png', '.jpg', '.jpeg'):
            return self._process_image(filepath, file_hash)
        elif ext == '.txt':
            return self._process_text(filepath, file_hash)
        return None

    def _process_pdf(self, filepath: Path, file_hash: str) -> ProcessedDocument:
        """Process PDF with pdfplumber — extracts text and tables."""
        all_text = []
        chunks = []

        with pdfplumber.open(filepath) as pdf:
            total_pages = len(pdf.pages)
            for i, page in enumerate(pdf.pages):
                # Extract text
                text = page.extract_text(layout=True)
                if text:
                    all_text.append(f"\n--- Page {i+1} ---\n{text}")

                # Extract tables
                tables = page.extract_tables()
                for t_idx, table in enumerate(tables):
                    if table and any(any(cell for cell in row) for row in table):
                        table_text = self._format_table(table)
                        all_text.append(f"\n[Table {t_idx+1} from Page {i+1}]\n{table_text}")

                # Create page-level chunk
                page_text = f"--- Page {i+1} ---\n"
                if text:
                    page_text += text
                if tables:
                    for t_idx, table in enumerate(tables):
                        if table:
                            page_text += f"\n[Table {t_idx+1}]:\n{self._format_table(table)}"

                if page_text.strip():
                    chunks.append(DocumentChunk(
                        text=page_text.strip(),
                        metadata={
                            "source": filepath.name,
                            "file_type": "pdf",
                            "page": i + 1,
                            "total_pages": total_pages,
                        },
                        chunk_id=f"{file_hash}_p{i+1}",
                        page_number=i + 1,
                    ))

        full_content = "\n".join(all_text)
        title = filepath.stem.replace('_', ' ').replace('-', ' ').title()

        return ProcessedDocument(
            filename=filepath.name,
            title=title,
            content=full_content,
            chunks=chunks,
            file_hash=file_hash,
            metadata={
                "source": filepath.name,
                "file_type": "pdf",
                "total_pages": total_pages,
                "total_chunks": len(chunks),
            }
        )

    def _process_pptx(self, filepath: Path, file_hash: str) -> ProcessedDocument:
        """Process PowerPoint presentations."""
        prs = Presentation(filepath)
        chunks = []
        all_text = []

        for slide_idx, slide in enumerate(prs.slides):
            slide_text_parts = []
            slide_text_parts.append(f"\n=== Slide {slide_idx + 1} ===")

            # Extract text from all shapes
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        para_text = paragraph.text.strip()
                        if para_text:
                            slide_text_parts.append(para_text)

                # Extract tables
                if shape.has_table:
                    table = shape.table
                    table_text = self._format_table_pptx(table)
                    slide_text_parts.append(f"\n[Table]:\n{table_text}")

            slide_text = "\n".join(slide_text_parts)
            if slide_text.strip():
                all_text.append(slide_text)
                chunks.append(DocumentChunk(
                    text=slide_text.strip(),
                    metadata={
                        "source": filepath.name,
                        "file_type": "pptx",
                        "slide": slide_idx + 1,
                        "total_slides": len(prs.slides),
                    },
                    chunk_id=f"{file_hash}_s{slide_idx + 1}",
                    page_number=slide_idx + 1,
                ))

        full_content = "\n".join(all_text)
        title = filepath.stem.replace('_', ' ').replace('-', ' ').title()

        return ProcessedDocument(
            filename=filepath.name,
            title=title,
            content=full_content,
            chunks=chunks,
            file_hash=file_hash,
            metadata={
                "source": filepath.name,
                "file_type": "pptx",
                "total_slides": len(prs.slides),
                "total_chunks": len(chunks),
            }
        )

    def _process_docx(self, filepath: Path, file_hash: str) -> ProcessedDocument:
        """Process Word documents."""
        try:
            from docx import Document as DocxDocument
        except ImportError:
            raise ImportError("python-docx not installed. Run: pip install python-docx")

        doc = DocxDocument(filepath)
        chunks = []
        all_text = []

        current_section = "Introduction"
        para_buffer = []
        chunk_idx = 0

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            # Detect headings
            if para.style.name.startswith('Heading'):
                if para_buffer:
                    chunk_text = "\n".join(para_buffer)
                    chunks.append(DocumentChunk(
                        text=chunk_text,
                        metadata={"source": filepath.name, "section": current_section},
                        chunk_id=f"{file_hash}_c{chunk_idx}",
                        section=current_section,
                    ))
                    chunk_idx += 1
                    para_buffer = []
                current_section = text
            else:
                para_buffer.append(text)

            all_text.append(text)

        # Flush remaining buffer
        if para_buffer:
            chunks.append(DocumentChunk(
                text="\n".join(para_buffer),
                metadata={"source": filepath.name, "section": current_section},
                chunk_id=f"{file_hash}_c{chunk_idx}",
                section=current_section,
            ))

        full_content = "\n".join(all_text)
        title = filepath.stem.replace('_', ' ').replace('-', ' ').title()

        return ProcessedDocument(
            filename=filepath.name,
            title=title,
            content=full_content,
            chunks=chunks,
            file_hash=file_hash,
            metadata={"source": filepath.name, "file_type": "docx", "total_chunks": len(chunks)}
        )

    def _process_image(self, filepath: Path, file_hash: str) -> Optional[ProcessedDocument]:
        """Process images with OCR."""
        try:
            image = Image.open(filepath)
            text = pytesseract.image_to_string(image)
            if not text.strip():
                return None

            chunks = [DocumentChunk(
                text=text.strip(),
                metadata={"source": filepath.name, "file_type": "image_ocr"},
                chunk_id=f"{file_hash}_img",
            )]

            return ProcessedDocument(
                filename=filepath.name,
                title=filepath.stem,
                content=text.strip(),
                chunks=chunks,
                file_hash=file_hash,
                metadata={"source": filepath.name, "file_type": "image_ocr"}
            )
        except Exception as e:
            print(f"  OCR failed for {filepath.name}: {e}")
            return None

    def _process_text(self, filepath: Path, file_hash: str) -> ProcessedDocument:
        """Process plain text files."""
        content = filepath.read_text(encoding='utf-8', errors='replace')
        lines = content.split('\n')

        chunks = []
        chunk_size = 500
        chunk_idx = 0

        for i in range(0, len(lines), chunk_size):
            chunk_text = "\n".join(lines[i:i + chunk_size])
            if chunk_text.strip():
                chunks.append(DocumentChunk(
                    text=chunk_text.strip(),
                    metadata={"source": filepath.name, "file_type": "txt"},
                    chunk_id=f"{file_hash}_t{chunk_idx}",
                ))
                chunk_idx += 1

        return ProcessedDocument(
            filename=filepath.name,
            title=filepath.stem,
            content=content,
            chunks=chunks,
            file_hash=file_hash,
            metadata={"source": filepath.name, "file_type": "txt", "total_chunks": len(chunks)}
        )

    def _format_table(self, table: list) -> str:
        """Format extracted table as readable text."""
        lines = []
        for row in table:
            if any(cell for cell in row):
                row_text = " | ".join(str(cell) if cell else "" for cell in row)
                lines.append(row_text)
        return "\n".join(lines)

    def _format_table_pptx(self, table) -> str:
        """Format PowerPoint table."""
        lines = []
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            lines.append(row_text)
        return "\n".join(lines)
